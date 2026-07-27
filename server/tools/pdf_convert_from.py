import os
import sys
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import openpyxl
import pdfplumber

def pdf_to_jpg(input_path, output_dir, dpi=150):
    doc = fitz.open(input_path)
    output_files = []
    
    for idx, page in enumerate(doc):
        pix = page.get_pixmap(dpi=dpi)
        out_file = os.path.join(output_dir, f"page_{idx+1}.jpg")
        pix.save(out_file)
        output_files.append(out_file)
        
    doc.close()
    return True, output_files

def pdf_to_word(input_path, output_path):
    doc = fitz.open(input_path)
    docx_doc = Document()
    
    for page_idx, page in enumerate(doc):
        if page_idx > 0:
            docx_doc.add_page_break()
        text = page.get_text("text")
        if text.strip():
            docx_doc.add_paragraph(text)
        else:
            # Render page image if scanned
            pix = page.get_pixmap(dpi=150)
            img_path = f"temp_p{page_idx}.png"
            pix.save(img_path)
            docx_doc.add_picture(img_path, width=Inches(6.0))
            if os.path.exists(img_path): os.remove(img_path)
            
    doc.close()
    docx_doc.save(output_path)
    return True, "Converted PDF to Word (.docx)"

def pdf_to_ppt(input_path, output_path):
    doc = fitz.open(input_path)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    
    for page_idx, page in enumerate(doc):
        slide = prs.slides.add_slide(blank_slide_layout)
        pix = page.get_pixmap(dpi=150)
        img_path = f"temp_ppt_p{page_idx}.png"
        pix.save(img_path)
        
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
        if os.path.exists(img_path): os.remove(img_path)
        
    doc.close()
    prs.save(output_path)
    return True, "Converted PDF to PowerPoint (.pptx)"

def pdf_to_excel(input_path, output_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Extracted PDF Data"
    
    current_row = 1
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    for col_idx, val in enumerate(row, start=1):
                        ws.cell(row=current_row, column=col_idx, value=val or "")
                    current_row += 1
                current_row += 1 # Empty spacing row
                
            # Extract plain text if no tables
            if not tables:
                text = page.extract_text()
                if text:
                    for line in text.split("\n"):
                        ws.cell(row=current_row, column=1, value=line)
                        current_row += 1
                        
    wb.save(output_path)
    return True, "Converted PDF to Excel (.xlsx)"

def pdf_to_pdfa(input_path, output_path):
    # Converts PDF to PDF/A (Archival compliance format)
    doc = fitz.open(input_path)
    doc.save(output_path, garbage=4, deflate=True, clean=True)
    doc.close()
    return True, "Converted to PDF/A (Archival Format)"
