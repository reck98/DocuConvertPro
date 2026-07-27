import unittest
import os
import sys
import shutil
import tempfile
import fitz # PyMuPDF
import pypdf
from reportlab.pdfgen import canvas
from docx import Document
from pptx import Presentation
import openpyxl

# Add server and server/tools to sys.path
SERVER_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TOOLS_DIR = os.path.join(SERVER_DIR, "tools")
sys.path.append(SERVER_DIR)
sys.path.append(TOOLS_DIR)

import pdf_organize
import pdf_optimize
import pdf_convert_to
import pdf_convert_from
import pdf_edit
import pdf_security
import pdf_ai
import pdf_workflow

class TestPdfToolsSuite(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.test_dir = tempfile.mkdtemp(prefix="pdf_test_")
        cls.sample_pdf_1 = os.path.join(cls.test_dir, "sample_1.pdf")
        cls.sample_pdf_2 = os.path.join(cls.test_dir, "sample_2.pdf")
        cls.sample_docx = os.path.join(cls.test_dir, "sample.docx")
        cls.sample_pptx = os.path.join(cls.test_dir, "sample.pptx")
        cls.sample_xlsx = os.path.join(cls.test_dir, "sample.xlsx")
        cls.sample_jpg = os.path.join(cls.test_dir, "sample.jpg")

        # 1. Create multi-page Sample PDF 1
        c1 = canvas.Canvas(cls.sample_pdf_1)
        c1.drawString(100, 750, "Page 1 - DocuConvert Test Document")
        c1.showPage()
        c1.drawString(100, 750, "Page 2 - Confidential Report SSN 123-45-6789")
        c1.showPage()
        c1.drawString(100, 750, "Page 3 - Financial Data Summary")
        c1.save()

        # 2. Create Sample PDF 2
        c2 = canvas.Canvas(cls.sample_pdf_2)
        c2.drawString(100, 750, "Page 1 - Appendix Section")
        c2.save()

        # 3. Create Sample DOCX
        doc = Document()
        doc.add_heading("DocuConvert Word Sample", 0)
        doc.add_paragraph("Testing Word to PDF conversion engine.")
        doc.save(cls.sample_docx)

        # 4. Create Sample PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "DocuConvert Presentation"
        prs.save(cls.sample_pptx)

        # 5. Create Sample XLSX
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sales"
        ws.cell(row=1, column=1, value="Product")
        ws.cell(row=1, column=2, value="Revenue")
        ws.cell(row=2, column=1, value="DocuConvert Pro")
        ws.cell(row=2, column=2, value="$50,000")
        wb.save(cls.sample_xlsx)

        # 6. Create Sample JPG
        pix = fitz.Pixmap(fitz.csRGB, fitz.Rect(0, 0, 100, 100), False)
        pix.clear_with(255)
        pix.save(cls.sample_jpg)

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls.test_dir, ignore_errors=True)

    # -------------------------------------------------------------------------
    # 1. ORGANIZE PDF TESTS
    # -------------------------------------------------------------------------
    def test_01_merge_pdfs(self):
        out_pdf = os.path.join(self.test_dir, "out_merge.pdf")
        success, msg = pdf_organize.merge_pdfs([self.sample_pdf_1, self.sample_pdf_2], out_pdf)
        self.assertTrue(success, f"Merge failed: {msg}")
        self.assertTrue(os.path.exists(out_pdf))
        doc = fitz.open(out_pdf)
        self.assertEqual(len(doc), 4) # 3 pages + 1 page = 4 pages
        doc.close()

    def test_02_split_pdf(self):
        success, files = pdf_organize.split_pdf(self.sample_pdf_1, self.test_dir, mode="all")
        self.assertTrue(success)
        self.assertGreaterEqual(len(files), 3)

    def test_03_remove_pages(self):
        out_pdf = os.path.join(self.test_dir, "out_remove.pdf")
        success, msg = pdf_organize.remove_pages(self.sample_pdf_1, out_pdf, "2")
        self.assertTrue(success)
        doc = fitz.open(out_pdf)
        self.assertEqual(len(doc), 2)
        doc.close()

    def test_04_extract_pages(self):
        out_pdf = os.path.join(self.test_dir, "out_extract.pdf")
        success, msg = pdf_organize.extract_pages(self.sample_pdf_1, out_pdf, "1,3")
        self.assertTrue(success)
        doc = fitz.open(out_pdf)
        self.assertEqual(len(doc), 2)
        doc.close()

    def test_05_organize_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_organize.pdf")
        success, msg = pdf_organize.organize_pdf(self.sample_pdf_1, out_pdf, page_orders=[3, 1, 2])
        self.assertTrue(success)
        self.assertTrue(os.path.exists(out_pdf))

    def test_06_scan_to_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_scan.pdf")
        success, msg = pdf_organize.scan_to_pdf(self.sample_pdf_1, out_pdf, grayscale=True, contrast=1.5)
        self.assertTrue(success)

    # -------------------------------------------------------------------------
    # 2. OPTIMIZE PDF TESTS
    # -------------------------------------------------------------------------
    def test_07_compress_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_compress.pdf")
        success, msg = pdf_optimize.compress_pdf(self.sample_pdf_1, out_pdf, level="medium")
        self.assertTrue(success)

    def test_08_repair_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_repair.pdf")
        success, msg = pdf_optimize.repair_pdf(self.sample_pdf_1, out_pdf)
        self.assertTrue(success)

    def test_09_ocr_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_ocr.pdf")
        success, msg = pdf_optimize.ocr_pdf(self.sample_pdf_1, out_pdf, lang="eng")
        self.assertTrue(success)

    # -------------------------------------------------------------------------
    # 3. CONVERT TO PDF TESTS
    # -------------------------------------------------------------------------
    def test_10_jpg_to_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_jpg2pdf.pdf")
        success, msg = pdf_convert_to.jpg_to_pdf([self.sample_jpg], out_pdf)
        self.assertTrue(success)

    def test_11_word_to_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_word2pdf.pdf")
        success, msg = pdf_convert_to.word_to_pdf(self.sample_docx, out_pdf)
        self.assertTrue(success)

    def test_12_ppt_to_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_ppt2pdf.pdf")
        success, msg = pdf_convert_to.ppt_to_pdf(self.sample_pptx, out_pdf)
        self.assertTrue(success)

    def test_13_excel_to_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_excel2pdf.pdf")
        success, msg = pdf_convert_to.excel_to_pdf(self.sample_xlsx, out_pdf)
        self.assertTrue(success)

    def test_14_html_to_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_html2pdf.pdf")
        html_code = "<h1>DocuConvert HTML Test</h1><p>Rendering HTML stream to PDF.</p>"
        success, msg = pdf_convert_to.html_to_pdf(html_code, out_pdf)
        self.assertTrue(success)

    # -------------------------------------------------------------------------
    # 4. CONVERT FROM PDF TESTS
    # -------------------------------------------------------------------------
    def test_15_pdf_to_jpg(self):
        success, files = pdf_convert_from.pdf_to_jpg(self.sample_pdf_1, self.test_dir, dpi=150)
        self.assertTrue(success)
        self.assertGreaterEqual(len(files), 1)

    def test_16_pdf_to_word(self):
        out_docx = os.path.join(self.test_dir, "out_pdf2word.docx")
        success, msg = pdf_convert_from.pdf_to_word(self.sample_pdf_1, out_docx)
        self.assertTrue(success)
        self.assertTrue(os.path.exists(out_docx))

    def test_17_pdf_to_ppt(self):
        out_pptx = os.path.join(self.test_dir, "out_pdf2ppt.pptx")
        success, msg = pdf_convert_from.pdf_to_ppt(self.sample_pdf_1, out_pptx)
        self.assertTrue(success)
        self.assertTrue(os.path.exists(out_pptx))

    def test_18_pdf_to_excel(self):
        out_xlsx = os.path.join(self.test_dir, "out_pdf2excel.xlsx")
        success, msg = pdf_convert_from.pdf_to_excel(self.sample_pdf_1, out_xlsx)
        self.assertTrue(success)
        self.assertTrue(os.path.exists(out_xlsx))

    def test_19_pdf_to_pdfa(self):
        out_pdfa = os.path.join(self.test_dir, "out_pdfa.pdf")
        success, msg = pdf_convert_from.pdf_to_pdfa(self.sample_pdf_1, out_pdfa)
        self.assertTrue(success)

    # -------------------------------------------------------------------------
    # 5. EDIT PDF TESTS
    # -------------------------------------------------------------------------
    def test_20_rotate_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_rotate.pdf")
        success, msg = pdf_edit.rotate_pdf(self.sample_pdf_1, out_pdf, angle=90)
        self.assertTrue(success)

    def test_21_add_page_numbers(self):
        out_pdf = os.path.join(self.test_dir, "out_numbers.pdf")
        success, msg = pdf_edit.add_page_numbers(self.sample_pdf_1, out_pdf, position="bottom-right")
        self.assertTrue(success)

    def test_22_add_watermark(self):
        out_pdf = os.path.join(self.test_dir, "out_watermark.pdf")
        success, msg = pdf_edit.add_watermark(self.sample_pdf_1, out_pdf, watermark_text="CONFIDENTIAL")
        self.assertTrue(success)

    def test_23_crop_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_crop.pdf")
        success, msg = pdf_edit.crop_pdf(self.sample_pdf_1, out_pdf, top=20, bottom=20)
        self.assertTrue(success)

    def test_24_edit_pdf_text(self):
        out_pdf = os.path.join(self.test_dir, "out_edittext.pdf")
        success, msg = pdf_edit.edit_pdf_text(self.sample_pdf_1, out_pdf)
        self.assertTrue(success)

    def test_25_fill_pdf_forms(self):
        out_pdf = os.path.join(self.test_dir, "out_forms.pdf")
        success, msg = pdf_edit.fill_pdf_forms(self.sample_pdf_1, out_pdf)
        self.assertTrue(success)

    # -------------------------------------------------------------------------
    # 6. PDF SECURITY TESTS
    # -------------------------------------------------------------------------
    def test_26_protect_and_unlock_pdf(self):
        out_protected = os.path.join(self.test_dir, "out_protected.pdf")
        success, msg = pdf_security.protect_pdf(self.sample_pdf_1, out_protected, user_password="PassWord123!")
        self.assertTrue(success)

        out_unlocked = os.path.join(self.test_dir, "out_unlocked.pdf")
        success, msg = pdf_security.unlock_pdf(out_protected, out_unlocked, password="PassWord123!")
        self.assertTrue(success)

    def test_27_sign_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_sign.pdf")
        success, msg = pdf_security.sign_pdf(self.sample_pdf_1, out_pdf, sign_text="Digitally Signed by Test")
        self.assertTrue(success)

    def test_28_redact_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_redact.pdf")
        success, msg = pdf_security.redact_pdf(self.sample_pdf_1, out_pdf, keywords=["Confidential", "SSN"])
        self.assertTrue(success)

    def test_29_compare_pdfs(self):
        out_pdf = os.path.join(self.test_dir, "out_compare.pdf")
        success, msg = pdf_security.compare_pdfs(self.sample_pdf_1, self.sample_pdf_2, out_pdf)
        self.assertTrue(success)

    # -------------------------------------------------------------------------
    # 7. AI INTELLIGENCE TESTS
    # -------------------------------------------------------------------------
    def test_30_ai_summarize_pdf(self):
        success, summary = pdf_ai.ai_summarize_pdf(self.sample_pdf_1)
        self.assertTrue(success)
        self.assertIn("total_words", summary)

    def test_31_translate_pdf(self):
        out_pdf = os.path.join(self.test_dir, "out_translate.pdf")
        success, msg = pdf_ai.translate_pdf(self.sample_pdf_1, out_pdf, target_lang="es")
        self.assertTrue(success)

    def test_32_pdf_to_markdown(self):
        out_md = os.path.join(self.test_dir, "out_markdown.md")
        success, content = pdf_ai.pdf_to_markdown(self.sample_pdf_1, out_md)
        self.assertTrue(success)
        self.assertTrue(os.path.exists(out_md))

    # -------------------------------------------------------------------------
    # 8. WORKFLOW AUTOMATION TEST
    # -------------------------------------------------------------------------
    def test_33_execute_workflow(self):
        steps = [
            {"action": "compress", "params": {"level": "low"}},
            {"action": "watermark", "params": {"text": "DRAFT"}}
        ]
        success, result = pdf_workflow.execute_workflow([self.sample_pdf_1], self.test_dir, steps)
        self.assertTrue(success)

if __name__ == "__main__":
    unittest.main()
